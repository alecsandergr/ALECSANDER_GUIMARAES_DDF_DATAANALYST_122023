import base64
from datetime import date
from io import BytesIO

import pandas as pd
import streamlit as st
from openai import OpenAI
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

from config import openai_key


# Function to get product details from ChatGPT API
def get_product_details(selected_product, df: pd.DataFrame):
    product_info = df[df["title"] == selected_product]
    if not product_info.empty:
        title = product_info.iloc[0]["title"]
        description = product_info.iloc[0]["text"]
        features = product_info.iloc[0]["features"]
        return {"title": title, "description": description, "features": features}
    else:
        return {"title": None, "description": None, "features": None}


def display_info(product_details: dict):
    # Display selected product
    st.subheader(f"Selected Product: {product_details['title']}")
    st.write(f"Description: {product_details['description']}")


def calculate_number_of_lines(text_frame, font_size, available_height):
    # Calculate the number of lines based on the font size
    first_paragraph = text_frame.paragraphs[0]
    line_height = font_size * (
        first_paragraph.line_spacing if first_paragraph.line_spacing else 1.25
    )

    return available_height // line_height


def check_font_size(text_frame, font_size, available_width, nb_of_lines):
    # Calculate the total width of the text in the TextFrame
    average_char_width = font_size / 2
    for paragraph in text_frame.paragraphs:
        total_width = 0
        for run in paragraph.runs:
            total_width += average_char_width * len(run.text)
        lines_spent = int(total_width / available_width) + (
            total_width % available_width > 0
        )
        nb_of_lines -= lines_spent

    if nb_of_lines >= 0:
        return False
    else:
        return True


def adjust_font_size(text_frame, content):
    original_text = text_frame.text

    # Calculate the available area based on the content dimensions
    available_width = content.width
    available_height = content.height

    # Set an initial font size
    font_size = Pt(24)
    nb_of_lines = calculate_number_of_lines(text_frame, font_size, available_height)

    # Loop to decrease font size until the text fits within the available area
    while (
        check_font_size(text_frame, font_size, available_width, nb_of_lines)
    ) and font_size > Pt(8):
        font_size = max(font_size - Pt(2), Pt(8))
        nb_of_lines = calculate_number_of_lines(text_frame, font_size, available_height)
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = font_size

    # If the font size has been adjusted, set the original text again
    if text_frame.text != original_text:
        text_frame.text = original_text


def generate_image(product_details: dict):
    client = OpenAI(api_key=openai_key)
    # prompt = f"""Build a image that best represents this discription
    # {product_details['description']} and title {product_details['title']}
    # """
    prompt = f"""Generate an image that visually represents the following product:
    Title: {product_details['title']}
    Description: {product_details['description']}
    Ensure that any text in the image is clear and readable.
    """
    response = client.images.generate(
        model="dall-e-3", prompt=prompt, size="1024x1024", response_format="b64_json"
    )

    return response.data[0].b64_json


def decode_image(b64_image_data):
    image_data = base64.b64decode(b64_image_data)
    image = Image.open(BytesIO(image_data))
    return image


def create_presentation(product_details: dict):
    prs = Presentation()

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]  # Title Slide Layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Data App"
    subtitle = slide.placeholders[1]
    subtitle.text = product_details["title"]

    # Slide 2: Product Description
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Description"
    content = slide.placeholders[1]
    content.text_frame.word_wrap = True
    content.text = product_details["description"]
    adjust_font_size(content.text_frame, content)

    # Slide 3: Features with Image
    slide_layout = prs.slide_layouts[1]  # Use the appropriate layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Features"
    # Calculate the center of the slide
    center_x = prs.slide_width / 2
    center_y = prs.slide_height / 2
    # Add the Word Cloud image
    left_img = Inches(0.5)
    top_img = center_y - Inches(2)
    height_img = Inches(5)
    width_img = Inches(4)
    b64_image_data = generate_image(product_details)
    image = decode_image(b64_image_data)
    img_bytes = BytesIO()
    image.save(img_bytes, format="PNG")
    slide.shapes.add_picture(
        img_bytes, left_img, top_img, width=width_img, height=height_img
    )
    # Add the product features with line breaks
    content = slide.placeholders[1]
    content.text_frame.word_wrap = True
    for feature, desc in product_details["features"].items():
        content.text += f"{feature}: {desc}\n"
    # Position the text on the right side
    left_text = center_x + Inches(0.5)  # Adjust left position for the text
    top_text = center_y - Inches(2)  # Adjust top position for the text
    content.left = round(left_text)
    content.top = round(top_text)
    content.width = Inches(4)  # Adjust width for the text
    content.height = Inches(5)
    adjust_font_size(content.text_frame, content)

    output_pptx = BytesIO()
    prs.save(output_pptx)
    prs.save("product_presentation.pptx")

    return output_pptx


def download_pptx(output_pptx: BytesIO, product_details: dict):
    # create file name
    name = "_".join(product_details["title"].split(" ")[:3])
    date_str = str(date.today()).replace("-", "_")
    filename = f"{name}_{date_str}.pptx"

    # display success message and download button
    st.success("The slides have been generated! :tada:")

    pptx = st.download_button(
        label="Click to download PowerPoint",
        data=output_pptx.getvalue(),
        file_name=filename,
    )

    return pptx


def get_download_link(outputx_pptx: BytesIO, product_details: dict) -> str:
    name = "_".join(product_details["title"].split(" ")[:3])
    date_str = str(date.today()).replace("-", "_")
    filename = f"{name}_{date_str}.pptx"

    # display success message and download button
    st.success("The slides have been generated! :tada:")
    # Convert the BytesIO object to base64
    b64 = base64.b64encode(outputx_pptx.getvalue()).decode()
    href = f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="{filename}">Download Presentation</a>'

    return href


def generate_presentation(df: pd.DataFrame):
    # Siderbar for product selection
    selected_product = st.sidebar.selectbox("Select a Product", df["title"])

    st.subheader("Product Presentation Generator")

    st.markdown(
        """
        Welcome to the Product Presentation Generator! Select a product from the sidebar to view its details
        and generate a PowerPoint presentation with main features, an image, and a word cloud.
        """
    )

    # Get the details of the selected product
    if st.sidebar.button("Generate Presentation"):
        product_details = get_product_details(selected_product, df)
        display_info(product_details)
        output_pptx = create_presentation(product_details)
        st.markdown(
            get_download_link(output_pptx, product_details), unsafe_allow_html=True
        )
